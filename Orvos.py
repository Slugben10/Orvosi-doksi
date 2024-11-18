import os
import requests
from flask import Flask, request, jsonify, render_template
from bs4 import BeautifulSoup
from langchain_community.document_loaders import TextLoader, PyPDFLoader
from langchain_community.embeddings import OpenAIEmbeddings
from langchain_community.vectorstores import Chroma
from langchain.text_splitter import CharacterTextSplitter
import openai
from docx import Document as DocxDocument
from langchain.schema import Document
import pandas as pd
from pptx import Presentation  # Library to load PowerPoint files

# OpenAI API key setup
os.environ['OPENAI_API_KEY'] = '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'

# Setup Flask app
app = Flask(__name__)

# OpenAI client configuration
openai.api_key = os.getenv('OPENAI_API_KEY')

# Initialize vector database
persist_directory = '/Users/binobenjamin/Documents/Orvosi-doksi/chroma_data'

# Function to load DOCX documents
def load_docx(file_path):
    doc = DocxDocument(file_path)
    return [Document(page_content=para.text, metadata={'source': file_path}) for para in doc.paragraphs if para.text]

# Function to load PPTX documents
def load_pptx(file_path):
    prs = Presentation(file_path)
    slides_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slides_text.append(shape.text)
    return [Document(page_content=text, metadata={'source': file_path}) for text in slides_text if text]

# Function to load XLSX documents
def load_xlsx(file_path):
    df = pd.read_excel(file_path, sheet_name=None)  # Load all sheets as dict of DataFrames
    documents = []
    for sheet_name, sheet_df in df.items():
        sheet_text = sheet_df.to_string(index=False, header=True)
        documents.append(Document(page_content=sheet_text, metadata={'source': file_path, 'sheet': sheet_name}))
    return documents

# Load and split documents from the specified directory
documents_folder = '/Users/binobenjamin/Documents/Orvosi-doksi/Dokumentumok'
all_documents = []

# Load documents from folder
for filename in os.listdir(documents_folder):
    if filename.endswith(".txt"):
        loader = TextLoader(os.path.join(documents_folder, filename))
        all_documents.extend(loader.load())
    elif filename.endswith(".docx"):
        docx_documents = load_docx(os.path.join(documents_folder, filename))
        all_documents.extend(docx_documents)
    elif filename.endswith(".pdf"):
        pdf_loader = PyPDFLoader(os.path.join(documents_folder, filename))
        all_documents.extend(pdf_loader.load())
    elif filename.endswith(".pptx"):
        pptx_documents = load_pptx(os.path.join(documents_folder, filename))
        all_documents.extend(pptx_documents)
    elif filename.endswith(".xlsx"):
        xlsx_documents = load_xlsx(os.path.join(documents_folder, filename))
        all_documents.extend(xlsx_documents)

# Check if documents are loaded correctly
# print("Loaded documents:", all_documents)

# Flatten the list to get only 'page_content'
documents_to_split = [doc for doc in all_documents if isinstance(doc, Document)]

# Split documents into chunks
text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=0)
documents = text_splitter.split_documents(documents_to_split)

# Store documents in vector database
db = Chroma.from_documents(
    documents,
    OpenAIEmbeddings(),
    persist_directory=persist_directory
)

# Web scraping function
def scrape_and_store(url):
    try:
        response = requests.get(url)
        if response.status_code == 200:
            soup = BeautifulSoup(response.content, 'html.parser')
            paragraphs = [p.get_text() for p in soup.find_all('p')]
            scraped_text = ' '.join(paragraphs)[:1000]

            # Split the scraped text into chunks
            scraped_documents = text_splitter.split_documents([Document(page_content=scraped_text)])
            
            # Store scraped documents in vector database
            db.add_documents(scraped_documents)
            return scraped_text
        else:
            return "Error: Unable to access the website."
    except Exception as e:
        return f"Error scraping the website: {str(e)}"

# Serve HTML frontend
@app.route('/')
def home():
    return render_template('index.html')

# Chat route to process user messages and call the GPT-3.5-turbo API
@app.route('/chat', methods=['POST'])
def chat():
    user_input = request.json.get('message')
    url_to_scrape = request.json.get('url')

    if not user_input:
        return jsonify({'response': "Error: No input provided"}), 400

    scraped_data = ''
    if url_to_scrape:
        scraped_data = scrape_and_store(url_to_scrape)

    try:
        # Search for matching document in vector database
        docs = db.similarity_search(user_input)
        if docs:
            document_match = docs[0].page_content  # Accessing 'page_content' from the Document object
        else:
            document_match = "No matching document found."

        # System message for API call
        system_message = (
            "You are a helpful assistant. Below is information scraped from a website and a related document that might be helpful:\n\n"
            f"Web scraped data: {scraped_data}\n\n"
            f"Relevant document: {document_match}\n\n"
            "Now, respond to the user's question in detail."
        )

        # Make an API call with the system message
        response = openai.ChatCompletion.create(
            model="gpt-4o-mini",  # Use a valid model name
            messages=[
                {"role": "system", "content": system_message},
                {"role": "user", "content": user_input},
            ]
        )
        bot_response = response['choices'][0]['message']['content']

        return jsonify({'response': bot_response})

    except Exception as e:
        return jsonify({'response': f"Error: {str(e)}"})

if __name__ == "__main__":
    app.run(debug=True)
